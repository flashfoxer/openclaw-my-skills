#!/usr/bin/env python3
"""
PPT to Video Converter - Main Pipeline Script
Converts PPTX to narrated video with AI voiceover and synced animations.

Usage:
    python ppt_to_video.py <pptx_path> [options]

Options:
    --voice        TTS voice (default: zh-CN-XiaoxiaoNeural)
    --rate         TTS speech rate (default: +10%)
    --output       Output video path (default: ./ppt-video-final.mp4)
    --zip          Also create ZIP package
    --subtitles    Add subtitles (default: off)
    --work-dir     Working directory for temp files

Requirements:
    pip install edge-tts mss Pillow lxml pywin32
    choco install ffmpeg
    WPS Office (for KWpp.Application COM)
"""

import os
import sys
import time
import subprocess
import ctypes
import json
import io
import re
import argparse
import tempfile
import zipfile
import shutil
import threading
from pathlib import Path
from PIL import Image
from threading import Thread
from queue import Queue

# ====== Defaults ======
DEFAULT_VOICE = "zh-CN-XiaoxiaoNeural"
DEFAULT_RATE = "+10%"
DEFAULT_OUTPUT = "ppt-video-final.mp4"
FFMPEG = r'C:\ProgramData\chocolatey\bin\ffmpeg.exe'
FFPROBE = r'C:\ProgramData\chocolatey\bin\ffprobe.exe'

user32 = ctypes.windll.user32


def log(msg):
    """Print with flush for real-time output"""
    print(msg, flush=True)


def find_wps_exe():
    """Find WPS Office executable path"""
    candidates = [
        r'C:\Program Files\WPS Office\office6\wpp.exe',
        r'C:\Program Files (x86)\WPS Office\office6\wpp.exe',
    ]
    for p in candidates:
        if os.path.exists(p):
            return p
    return None


def parse_ppt(pptx_path):
    """
    Extract slide count, text content, and animation counts from PPTX.
    Returns (slide_count, slide_texts, anim_counts)
    """
    from pptx import Presentation
    from lxml import etree
    import zipfile

    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)

    slide_texts = []
    anim_counts = []

    # Extract animation counts from XML
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        slide_files = sorted(
            [f for f in zf.namelist()
             if f.startswith('ppt/slides/slide') and f.endswith('.xml')],
            key=lambda x: int(re.search(r'slide(\d+)', x).group(1))
        )

        for slide_file in slide_files:
            content = zf.read(slide_file)
            root = etree.fromstring(content)
            count = 0
            for elem in root.iter():
                tag = etree.QName(elem.tag).localname
                if tag == 'cond':
                    delay = elem.get('delay', '')
                    if delay == 'indefinite':
                        count += 1
            anim_counts.append(count)

    # Extract text from each slide
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(' | '.join(row_texts))
        slide_texts.append(texts)

    return slide_count, slide_texts, anim_counts


def generate_tts_scripts(slide_texts, slide_count):
    """
    Generate narration scripts for each slide based on extracted text.
    Override this function to customize narration.
    """
    scripts = {}
    for i in range(slide_count):
        texts = slide_texts[i]
        if texts:
            # Join all text from the slide into narration
            script = '。'.join(texts[:3])  # Use first 3 text blocks
            if len(script) > 200:
                script = script[:200] + '...'
            scripts[i + 1] = script
        else:
            scripts[i + 1] = f"第{i+1}页。"
    return scripts


async def generate_tts_audio(scripts, audio_dir, voice=DEFAULT_VOICE, rate=DEFAULT_RATE):
    """Generate TTS audio files for each slide"""
    import edge_tts

    os.makedirs(audio_dir, exist_ok=True)

    for i, (slide_num, text) in enumerate(sorted(scripts.items()), 1):
        out_path = os.path.join(audio_dir, f'slide{slide_num:02d}.mp3')
        log(f"  Slide {slide_num:02d}: {len(text)}字 ...")

        communicate = edge_tts.Communicate(text, voice, rate=rate)
        await communicate.save(out_path)

    log(f"  TTS生成完成!")


def measure_audio_durations(audio_dir, slide_count):
    """Measure duration of each audio file"""
    durations = []
    for i in range(1, slide_count + 1):
        path = os.path.join(audio_dir, f'slide{i:02d}.mp3')
        result = subprocess.run(
            [FFPROBE, '-v', 'error', '-show_entries', 'format=duration',
             '-of', 'default=noprint_wrappers=1:nokey=1', path],
            capture_output=True, text=True
        )
        dur = float(result.stdout.strip())
        durations.append(dur)
    return durations


def fix_ppt_animations(src_pptx, dst_pptx):
    """
    Modify PPTX to auto-play all animations (no click required).
    Changes delay="indefinite" to delay="0" in slide XML.
    """
    import zipfile
    from lxml import etree
    import tempfile

    tmpdir = tempfile.mkdtemp()

    with zipfile.ZipFile(src_pptx, 'r') as zf:
        zf.extractall(tmpdir)

    slides_dir = os.path.join(tmpdir, 'ppt', 'slides')
    if not os.path.exists(slides_dir):
        shutil.rmtree(tmpdir)
        shutil.copy2(src_pptx, dst_pptx)
        return 0

    slide_files = sorted([f for f in os.listdir(slides_dir) if f.endswith('.xml')])
    total_changes = 0

    for slide_file in slide_files:
        slide_path = os.path.join(slides_dir, slide_file)
        tree = etree.parse(slide_path)
        root = tree.getroot()

        changed = 0
        for elem in root.iter():
            tag = etree.QName(elem.tag).localname
            if tag == 'cond':
                delay = elem.get('delay', '')
                if delay == 'indefinite':
                    elem.set('delay', '0')
                    changed += 1

        if changed > 0:
            tree.write(slide_path, xml_declaration=True, encoding='UTF-8', standalone=True)
            total_changes += changed

    # Re-zip
    with zipfile.ZipFile(dst_pptx, 'w', zipfile.ZIP_DEFLATED) as zf:
        for root_dir, dirs, files in os.walk(tmpdir):
            for file in files:
                file_path = os.path.join(root_dir, file)
                arcname = os.path.relpath(file_path, tmpdir)
                zf.write(file_path, arcname)

    shutil.rmtree(tmpdir)
    return total_changes


def send_key(vk, hold=0.05):
    """Send a keyboard event"""
    user32.keybd_event(vk, 0, 0, 0)
    time.sleep(hold)
    user32.keybd_event(vk, 0, 0x0002, 0)
    time.sleep(hold)


def launch_wps_slideshow(pptx_path):
    """Launch WPS slideshow via PowerShell COM and return the process"""
    ps_script = f'''
try {{
    $wpp = New-Object -ComObject KWpp.Application
    $wpp.Visible = $true
    $pres = $wpp.Presentations.Open("{pptx_path}")
    $pres.SlideShowSettings.Run()
    while ($true) {{ Start-Sleep 5 }}
}} catch {{
    Write-Host "Error: $_"
}}
'''

    ps_proc = subprocess.Popen(
        ['powershell', '-Command', ps_script],
        stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True
    )

    def read_ps():
        while ps_proc.poll() is None:
            line = ps_proc.stdout.readline()
            if line:
                log(f"  {line.strip()}")

    threading.Thread(target=read_ps, daemon=True).start()
    return ps_proc


def record_slideshow(pptx_path, total_dur, slide_durations, frames_dir):
    """
    Record PPT slideshow with screen capture.
    Animations auto-play, slides advance via right arrow at slide_durations boundaries.
    """
    import mss

    # Clean frames dir
    os.makedirs(frames_dir, exist_ok=True)
    for f in os.listdir(frames_dir):
        os.remove(os.path.join(frames_dir, f))

    # Calculate slide end times
    slide_ends = []
    cum = 0
    for d in slide_durations:
        cum += d
        slide_ends.append(cum)

    SW, SH = user32.GetSystemMetrics(0), user32.GetSystemMetrics(1)
    log(f"  屏幕: {SW}x{SH}")

    # Kill old processes
    subprocess.run(['taskkill', '/f', '/im', 'wpp.exe'], capture_output=True)
    subprocess.run(['taskkill', '/f', '/im', 'wps.exe'], capture_output=True)
    time.sleep(3)

    # Launch WPS
    log("  启动WPS放映...")
    ps_proc = launch_wps_slideshow(pptx_path)
    time.sleep(15)
    log("  WPS放映已启动!")

    # Writer thread
    frame_queue = Queue(maxsize=3000)
    writer_done = [False]

    def writer():
        while not writer_done[0] or not frame_queue.empty():
            try:
                idx, buf = frame_queue.get(timeout=3)
                with open(os.path.join(frames_dir, f"frame_{idx:08d}.jpg"), 'wb') as f:
                    f.write(buf.getvalue())
                frame_queue.task_done()
            except:
                pass

    Thread(target=writer, daemon=True).start()

    # Capture + Control
    log(f"  录制中 ({total_dur:.1f}s)...")
    user32.ShowCursor(False)
    start_time = time.time()
    frame_idx = 0
    slide_idx = 0
    next_report = 0

    with mss.mss() as sct:
        monitor = sct.monitors[1]

        while True:
            elapsed = time.time() - start_time
            if elapsed >= total_dur:
                break

            img = sct.grab(monitor)
            pil = Image.frombytes('RGB', img.size, img.bgra, 'raw', 'BGRX')
            buf = io.BytesIO()
            pil.save(buf, format='JPEG', quality=92)
            frame_queue.put((frame_idx, buf))
            frame_idx += 1

            # Advance slide when audio is done
            while slide_idx < len(slide_ends) - 1 and elapsed >= slide_ends[slide_idx] - 0.1:
                log(f"    >>> 翻页: {slide_idx+1} -> {slide_idx+2} (t={elapsed:.1f}s)")
                send_key(0x27, 0.15)  # RIGHT ARROW
                slide_idx += 1
                time.sleep(0.3)

            if elapsed >= next_report:
                fps_eff = frame_idx / max(elapsed, 0.01)
                log(f"   {elapsed:.0f}/{total_dur:.0f}s | {frame_idx}帧 | ~{fps_eff:.0f}fps | 幻灯{slide_idx+1}/{len(slide_durations)}")
                next_report = elapsed + 15

    real_time = time.time() - start_time
    writer_done[0] = True
    time.sleep(8)

    # Exit
    send_key(0x1B, 0.2)  # ESC
    time.sleep(2)
    subprocess.run(['taskkill', '/f', '/im', 'wpp.exe'], capture_output=True)
    subprocess.run(['taskkill', '/f', '/im', 'wps.exe'], capture_output=True)
    ps_proc.kill()
    user32.ShowCursor(True)

    total_frames = frame_idx

    # Remove invalid frames
    valid_count = 0
    for f in os.listdir(frames_dir):
        path = os.path.join(frames_dir, f)
        if os.path.getsize(path) < 1000:
            os.remove(path)
        else:
            valid_count += 1

    calc_fps = valid_count / total_dur
    log(f"  录制完成: {valid_count}有效帧, fps={calc_fps:.2f}")
    return valid_count, calc_fps


def encode_video_ffmpeg(frames_dir, calc_fps, output_path):
    """Encode frames to video using FFmpeg"""
    result = subprocess.run([
        FFMPEG, '-y',
        '-framerate', str(calc_fps),
        '-i', os.path.join(frames_dir, 'frame_%08d.jpg'),
        '-c:v', 'libx264', '-pix_fmt', 'yuv420p',
        '-preset', 'fast', '-crf', '23', '-r', '30',
        output_path
    ], capture_output=True, text=True, timeout=600)

    if result.returncode != 0:
        log(f"  FFmpeg错误: {result.stderr[-500:]}")
        return False
    return True


def encode_video_opencv(frames_dir, calc_fps, output_path):
    """Fallback: encode frames to video using OpenCV"""
    import cv2

    frame_files = sorted([f for f in os.listdir(frames_dir) if f.endswith('.jpg')])
    if not frame_files:
        return False

    img = cv2.imread(os.path.join(frames_dir, frame_files[0]))
    h, w = img.shape[:2]

    fourcc = cv2.VideoWriter_fourcc(*'XVID')
    out = cv2.VideoWriter(output_path, fourcc, calc_fps, (w, h), True)

    for fname in frame_files:
        frame = cv2.imread(os.path.join(frames_dir, fname))
        if frame is not None:
            if len(frame.shape) == 2:
                frame = cv2.cvtColor(frame, cv2.COLOR_GRAY2BGR)
            out.write(frame)

    out.release()
    return True


def merge_audio(audio_dir, slide_count, output_path):
    """Concatenate per-slide audio files"""
    with open(output_path, 'wb') as out_f:
        for i in range(1, slide_count + 1):
            path = os.path.join(audio_dir, f'slide{i:02d}.mp3')
            if os.path.exists(path):
                with open(path, 'rb') as in_f:
                    out_f.write(in_f.read())


def final_merge(video_path, audio_path, output_path, subtitles_srt=None):
    """Merge video and audio (optionally with subtitles) into final MP4"""
    cmd = [
        FFMPEG, '-y',
        '-i', video_path, '-i', audio_path,
        '-c:v', 'libx264', '-pix_fmt', 'yuv420p', '-preset', 'fast', '-crf', '23',
        '-c:a', 'aac', '-b:a', '128k', '-shortest',
    ]

    if subtitles_srt and os.path.exists(subtitles_srt):
        srt_filter = subtitles_srt.replace('\\', '/').replace(':', '\\:')
        style = ("FontName=SimSun,FontSize=12,PrimaryColour=&H00FFFFFF,"
                 "OutlineColour=&H00000000,BackColour=&H66000000,"
                 "Bold=0,Italic=0,Alignment=2,MarginV=30,BorderStyle=4,Outline=0,Shadow=0")
        cmd.extend(['-vf', f"subtitles='{srt_filter}':force_style='{style}'"])

    cmd.append(output_path)

    result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
    return result.returncode == 0


def generate_srt(scripts, slide_durations, output_path, max_chars=24):
    """Generate SRT subtitle file"""
    def split_text(text, mc=max_chars):
        ts = re.split(r'([，,。！？、；：])', text)
        chunks, cur = [], ""
        i = 0
        while i < len(ts):
            cur += ts[i]
            if i + 1 < len(ts) and ts[i + 1] in '，,。！？、；：':
                cur += ts[i + 1]
                i += 1
            i += 1
            if len(cur) > mc or i >= len(ts):
                if cur.strip():
                    chunks.append(cur.strip())
                cur = ""
        return chunks

    def fmt_time(t):
        h, m, s = int(t // 3600), int((t % 3600) // 60), int(t % 60)
        return f"{h:02d}:{m:02d}:{s:02d},{int((t - int(t)) * 1000):03d}"

    cum_t, idx, lines = 0.0, 1, []
    for i in range(1, len(scripts) + 1):
        chunks = split_text(scripts[i])
        dur = slide_durations[i - 1]
        if len(chunks) <= 1:
            lines.append(f"{idx}\n{fmt_time(cum_t)} --> {fmt_time(cum_t + dur)}\n{chunks[0]}\n")
            idx += 1
        else:
            gap = 0.05
            tgap = gap * (len(chunks) - 1)
            per = (dur - tgap) / len(chunks)
            for ci, c in enumerate(chunks):
                s = cum_t + ci * (per + gap)
                e = s + per
                if ci == len(chunks) - 1:
                    e = cum_t + dur
                lines.append(f"{idx}\n{fmt_time(s)} --> {fmt_time(e)}\n{c}\n")
                idx += 1
        cum_t += dur

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))


def binary_copy(src, dst):
    """Binary copy file (handles permission issues)"""
    with open(src, 'rb') as f:
        data = f.read()
    with open(dst, 'wb') as f:
        f.write(data)
    return os.path.getsize(dst)


def main():
    parser = argparse.ArgumentParser(description='Convert PPT to video with AI voiceover')
    parser.add_argument('pptx', help='Path to PPTX file')
    parser.add_argument('--voice', default=DEFAULT_VOICE, help=f'TTS voice (default: {DEFAULT_VOICE})')
    parser.add_argument('--rate', default=DEFAULT_RATE, help=f'TTS rate (default: {DEFAULT_RATE})')
    parser.add_argument('--output', default=DEFAULT_OUTPUT, help='Output video path')
    parser.add_argument('--zip', action='store_true', help='Also create ZIP package')
    parser.add_argument('--subtitles', action='store_true', help='Add subtitles')
    parser.add_argument('--work-dir', default=None, help='Working directory for temp files')
    args = parser.parse_args()

    pptx_path = os.path.abspath(args.pptx)
    if not os.path.exists(pptx_path):
        log(f"Error: PPTX not found: {pptx_path}")
        sys.exit(1)

    # Setup work directory
    work_dir = args.work_dir or os.path.join(tempfile.gettempdir(), 'ppt_to_video')
    audio_dir = os.path.join(work_dir, 'audio')
    frames_dir = os.path.join(work_dir, 'frames')
    fixed_pptx = os.path.join(work_dir, 'presentation_auto.pptx')

    os.makedirs(audio_dir, exist_ok=True)
    os.makedirs(frames_dir, exist_ok=True)

    log("=" * 60)
    log("PPT to Video Converter")
    log("=" * 60)

    # Step 1: Parse PPT
    log("\n[1/7] 解析PPT...")
    slide_count, slide_texts, anim_counts = parse_ppt(pptx_path)
    log(f"  幻灯片: {slide_count} | 动画: {sum(anim_counts)}")

    # Step 2: Generate TTS
    log("\n[2/7] 生成TTS配音...")
    scripts = generate_tts_scripts(slide_texts, slide_count)
    import asyncio
    asyncio.run(generate_tts_audio(scripts, audio_dir, args.voice, args.rate))

    # Measure durations
    slide_durations = measure_audio_durations(audio_dir, slide_count)
    total_dur = sum(slide_durations)
    log(f"  总时长: {total_dur:.1f}s ({total_dur/60:.1f}min)")

    # Save durations
    with open(os.path.join(work_dir, 'durations.json'), 'w') as f:
        json.dump(slide_durations, f)

    # Step 3: Fix animations
    log("\n[3/7] 修改PPT动画为自动播放...")
    changes = fix_ppt_animations(pptx_path, fixed_pptx)
    log(f"  修改: {changes}个动画")

    # Step 4: Record
    log("\n[4/7] 录制PPT放映...")
    valid_frames, calc_fps = record_slideshow(
        fixed_pptx, total_dur, slide_durations, frames_dir
    )

    # Step 5: Encode
    log(f"\n[5/7] 编码视频 ({calc_fps:.2f}fps)...")
    temp_video = os.path.join(work_dir, 'video.mp4')
    if not encode_video_ffmpeg(frames_dir, calc_fps, temp_video):
        log("  FFmpeg失败, 尝试OpenCV...")
        temp_video = os.path.join(work_dir, 'video.avi')
        encode_video_opencv(frames_dir, calc_fps, temp_video)

    # Step 6: Merge audio
    log("\n[6/7] 合并音频...")
    merged_audio = os.path.join(work_dir, 'merged.mp3')
    merge_audio(audio_dir, slide_count, merged_audio)

    # Generate subtitles if requested
    srt_path = None
    if args.subtitles:
        srt_path = os.path.join(work_dir, 'subs.srt')
        generate_srt(scripts, slide_durations, srt_path)

    # Step 7: Final merge
    log("\n[7/7] 最终合成...")
    temp_final = os.path.join(work_dir, 'final.mp4')
    if not final_merge(temp_video, merged_audio, temp_final, srt_path):
        log("  合成失败!")
        sys.exit(1)

    # Copy to output
    output_path = os.path.abspath(args.output)
    size = binary_copy(temp_final, output_path)
    log(f"  视频: {size/1024/1024:.2f}MB")

    # ZIP if requested
    if args.zip:
        zip_path = output_path.replace('.mp4', '.zip')
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.write(output_path, os.path.basename(output_path))
        log(f"  ZIP: {os.path.getsize(zip_path)/1024/1024:.2f}MB")

    # Verify
    probe = subprocess.run(
        [FFPROBE, '-v', 'error',
         '-show_entries', 'stream=width,height,r_frame_rate,duration',
         '-of', 'default=noprint_wrappers=1', output_path],
        capture_output=True, text=True
    )
    log(f"\n验证:\n{probe.stdout.strip()}")
    log("\n完成!")


if __name__ == '__main__':
    main()